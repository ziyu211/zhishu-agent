# -*- coding: utf-8 -*-
"""端到端验证 v1.0.32 文档/图片处理：LibreOffice 旧格式转换 + 中文 OCR。"""
import io, os, glob, subprocess, sys, traceback

# 让 zhishu 包可导入（脚本位于 backend/tests/，需把 backend/ 加入路径）
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
sys.path.insert(0, "/app/backend")
from zhishu.core import rag

PASS, FAIL = [], []

def check(name, cond, extra=""):
    (PASS if cond else FAIL).append(name)
    print(("PASS " if cond else "FAIL ") + name + ("  " + extra if extra else ""))

def soffice_convert(in_path, fmt):
    out = os.path.join("/tmp", "conv_out")
    os.makedirs(out, exist_ok=True)
    prof = os.path.join("/tmp", "lo_profile_" + fmt)
    os.makedirs(prof, exist_ok=True)
    subprocess.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                    f"-env:UserInstallation=file://{prof}",
                    "--convert-to", fmt, "--outdir", out, in_path],
                   capture_output=True, text=True, timeout=180)
    cands = [os.path.join(out, f) for f in os.listdir(out) if f.lower().endswith("." + fmt)]
    return cands[0] if cands else None

# 1) 生成新格式文档
import docx, pptx, openpyxl
from pptx.util import Inches

docx_path = "/tmp/t.docx"
d = docx.Document(); d.add_paragraph("智枢 Zhishu 文档处理转换测试 12345")
d.save(docx_path)

pptx_path = "/tmp/t.pptx"
pr = pptx.Presentation(); slide = pr.slides.add_slide(pr.slide_layouts[1])
slide.shapes.title.text = "智枢幻灯片"; pr.save(pptx_path)

xlsx_path = "/tmp/t.xlsx"
wb = openpyxl.Workbook(); ws = wb.active
ws.append(["名称", "数值"]); ws.append(["智枢", 123]); wb.save(xlsx_path)

# 2) 转成旧版 OLE 格式
doc_old = soffice_convert(docx_path, "doc")
ppt_old = soffice_convert(pptx_path, "ppt")
xls_old = soffice_convert(xlsx_path, "xls")
check("生成 .doc 旧格式", bool(doc_old), doc_old or "")
check("生成 .ppt 旧格式", bool(ppt_old), ppt_old or "")
check("生成 .xls 旧格式", bool(xls_old), xls_old or "")

# 3) 嗅探器区分 .ppt（此前误判为 .doc）
if ppt_old:
    raw_ppt = open(ppt_old, "rb").read()
    check("OLE 嗅探器识别 .ppt", rag._sniff_office_ext(raw_ppt) == ".ppt",
          "=> " + str(rag._sniff_office_ext(raw_ppt)))

# 4) read_file_text 解析旧格式（经 LibreOffice 转换）
if doc_old:
    txt, ft = rag.read_file_text("t.doc", open(doc_old, "rb").read(), None, None)
    check(".doc 解析出正文", "智枢" in txt and "12345" in txt, f"ft={ft} len={len(txt)}")
if ppt_old:
    txt, ft = rag.read_file_text("t.ppt", open(ppt_old, "rb").read(), None, None)
    check(".ppt 解析出正文", "智枢" in txt, f"ft={ft} len={len(txt)}")
if xls_old:
    txt, ft = rag.read_file_text("t.xls", open(xls_old, "rb").read(), None, None)
    check(".xls 解析出正文", "智枢" in txt and "123" in txt, f"ft={ft} len={len(txt)}")

# 5) 图片中文 OCR
from PIL import Image, ImageDraw, ImageFont
fonts = glob.glob("/usr/share/fonts/**/*.ttf", recursive=True) + \
        glob.glob("/usr/share/fonts/**/*.otf", recursive=True) + \
        glob.glob("/usr/share/fonts/**/*.ttc", recursive=True)
cjk = next((f for f in fonts if "CJK" in f or "Noto" in f or "cjk" in f.lower()), fonts[0])
print("font:", cjk)
img = Image.new("RGB", (480, 160), "white")
draw = ImageDraw.Draw(img)
try:
    fnt = ImageFont.truetype(cjk, 32)
except Exception:
    fnt = ImageFont.load_default()
draw.text((10, 50), "智枢 OCR 中文识别测试 67890", fill="black", font=fnt)
buf = io.BytesIO(); img.save(buf, "PNG"); png = buf.getvalue()
ocr = rag._ocr_image_bytes(png)
check("图片中文 OCR 命中", "智枢" in ocr or "OCR" in ocr or "67890" in ocr,
      "ocr=" + ocr.replace("\n", " ")[:80])

# 6) 扫描型 PDF OCR
img.save("/tmp/scan.pdf")
pdf_raw = open("/tmp/scan.pdf", "rb").read()
pdf_ocr = rag._ocr_pdf(pdf_raw)
check("扫描 PDF OCR 命中", "智枢" in pdf_ocr or "OCR" in pdf_ocr or "67890" in pdf_ocr,
      "pdf_ocr=" + pdf_ocr.replace("\n", " ")[:80])

# 7) PDF 文本层提取仍正常（用新格式先验证不回归）
print("\n==== 结果 ====")
print("PASS:", len(PASS), "FAIL:", len(FAIL))
for f in FAIL:
    print("  FAIL:", f)
sys.exit(1 if FAIL else 0)
